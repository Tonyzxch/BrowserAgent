#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
浏览器智能助手 - 项目介绍PPT生成器
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 设置主题颜色
    PRIMARY_COLOR = RGBColor(102, 126, 234)  # 紫色 #667eea
    SECONDARY_COLOR = RGBColor(118, 75, 162)  # 深紫色 #764ba2
    TEXT_COLOR = RGBColor(30, 41, 59)  # 深灰色
    LIGHT_GRAY = RGBColor(241, 245, 249)

    def add_title_slide(title, subtitle):
        """添加标题幻灯片"""
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 渐变背景（用矩形模拟）
        left = Inches(0)
        top = Inches(0)
        width = prs.slide_width
        height = prs.slide_height
        shape = slide.shapes.add_shape(1, left, top, width, height)
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = PRIMARY_COLOR
        shape.line.fill.background()

        # 标题
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(1.5)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 副标题
        left = Inches(1)
        top = Inches(4.2)
        width = Inches(8)
        height = Inches(1)
        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        return slide

    def add_content_slide(title, title_en, content_items):
        """添加内容幻灯片"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 背景
        left = Inches(0)
        top = Inches(0)
        width = prs.slide_width
        height = prs.slide_height
        shape = slide.shapes.add_shape(1, left, top, width, height)
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.line.fill.background()

        # 标题栏（带渐变效果）
        left = Inches(0)
        top = Inches(0)
        width = prs.slide_width
        height = Inches(1.2)
        title_bg = slide.shapes.add_shape(1, left, top, width, height)
        fill = title_bg.fill
        fill.solid()
        fill.fore_color.rgb = PRIMARY_COLOR
        title_bg.line.fill.background()

        # 中文标题
        left = Inches(0.5)
        top = Inches(0.25)
        width = Inches(9)
        height = Inches(0.5)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # 英文标题
        left = Inches(0.5)
        top = Inches(0.7)
        width = Inches(9)
        height = Inches(0.3)
        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = title_en
        subtitle_frame.paragraphs[0].font.size = Pt(18)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        subtitle_frame.paragraphs[0].font.italic = True

        # 内容区域
        left = Inches(0.8)
        top = Inches(1.8)
        width = Inches(8.4)
        height = Inches(5)

        content_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for item in content_items:
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(20)
            p.font.color.rgb = TEXT_COLOR
            p.space_before = Pt(12)
            p.space_after = Pt(12)

        return slide

    # 第1页：封面
    add_title_slide(
        "浏览器智能助手\nBrowser Agent Assistant",
        "AI驱动的自主浏览器操作扩展"
    )

    # 第2页：项目背景
    add_content_slide(
        "项目背景",
        "Background",
        [
            "🎯 任务目标：开发智能浏览器插件，通过自然语言实现网页自动化",
            "",
            "🔍 核心挑战：",
            "  • 理解用户意图和任务目标",
            "  • 自动感知页面结构和可交互元素",
            "  • 智能规划多步骤执行流程",
            "  • 准确定位元素并执行操作",
            "",
            "💡 解决方案：AI Agent + Chrome Extension",
        ]
    )

    # 第3页：核心功能
    add_content_slide(
        "核心功能",
        "Key Features",
        [
            "🤖 自然语言交互",
            "   用中文描述任务，AI自动理解并执行",
            "",
            "🔍 智能页面感知",
            "   自动识别可交互元素，生成精准选择器",
            "",
            "🧠 AI任务规划",
            "   Claude Opus 4.8驱动的多步骤任务分解",
            "",
            "⚡ 自动化执行",
            "   支持点击、输入、导航、信息提取等操作",
        ]
    )

    # 第4页：技术架构
    add_content_slide(
        "技术架构",
        "Technical Architecture",
        [
            "📱 前端技术",
            "   React 18 + TypeScript + Vite 5",
            "",
            "🔌 Chrome Extension",
            "   Manifest V3 | Background Service Worker | Content Script",
            "",
            "🤖 AI服务",
            "   Claude Opus 4.8 | OpenAI兼容API",
            "",
            "🐳 部署方案",
            "   Docker + Nginx | 一键部署",
        ]
    )

    # 第5页：设计亮点
    add_content_slide(
        "设计亮点",
        "Design Highlights",
        [
            "🎯 智能元素定位",
            "   CSS Selector + XPath双重策略，自动过滤隐藏元素",
            "",
            "🧠 结构化AI交互",
            "   精心设计的Prompt工程，JSON格式化响应",
            "",
            "📊 实时状态同步",
            "   Chrome Messaging API实现多组件响应式更新",
            "",
            "🏛️ 模块化架构",
            "   清晰的职责分离，TypeScript类型安全",
        ]
    )

    # 第6页：功能演示
    add_content_slide(
        "使用示例",
        "Use Cases",
        [
            "📌 场景1：GitHub项目搜索",
            "   \"帮我在GitHub搜索React UI组件库，按star排序\"",
            "   → 自动定位搜索框 → 输入关键词 → 应用筛选",
            "",
            "📌 场景2：豆瓣电影查询",
            "   \"在豆瓣搜索评分8分以上的科幻电影\"",
            "   → 导航到豆瓣 → 输入搜索 → 设置筛选 → 提取信息",
            "",
            "📌 场景3：电商商品筛选",
            "   \"在淘宝搜索500元以内的蓝牙耳机\"",
            "   → 搜索商品 → 设置价格区间 → 提取列表",
        ]
    )

    # 第7页：实现成果
    add_content_slide(
        "实现成果",
        "Achievements",
        [
            "✅ 功能完整性",
            "   任务理解、页面感知、规划执行、结果反馈完整闭环",
            "",
            "✅ 代码质量",
            "   2000+行代码，TypeScript类型安全，模块化设计",
            "",
            "✅ 文档完整",
            "   11个Markdown文档，涵盖安装、使用、开发全流程",
            "",
            "✅ 生产就绪",
            "   Docker一键部署，插件大小仅65KB",
        ]
    )

    # 第8页：开发过程
    add_content_slide(
        "开发过程",
        "Development Process",
        [
            "⏱️ 开发时长：2小时（使用Claude Code）",
            "",
            "📝 迭代过程：",
            "   1️⃣ 需求分析与架构设计（10分钟）",
            "   2️⃣ 核心功能开发（60分钟）",
            "   3️⃣ UI界面与交互优化（30分钟）",
            "   4️⃣ Docker部署配置（15分钟）",
            "   5️⃣ 文档编写与完善（15分钟）",
            "",
            "🛠️ 开发工具：Claude Code + Claude Opus 4.7",
        ]
    )

    # 第9页：未来展望
    add_content_slide(
        "未来规划",
        "Future Plans",
        [
            "🚀 功能增强",
            "   • 支持更多动作类型（拖拽、右键菜单）",
            "   • 添加任务模板系统",
            "   • 多标签页协同操作",
            "",
            "⚡ 性能优化",
            "   • 减少页面信息传输量",
            "   • 缓存AI规划结果",
            "   • 并行执行独立步骤",
            "",
            "🌐 平台扩展",
            "   • 支持Firefox等浏览器",
            "   • 添加本地模型选项",
        ]
    )

    # 第10页：联系方式
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    left = Inches(0)
    top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height
    shape = slide.shapes.add_shape(1, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()

    # 感谢标题
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "Thank You!"
    title_frame.paragraphs[0].font.size = Pt(60)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 项目信息
    left = Inches(1)
    top = Inches(3.5)
    width = Inches(8)
    height = Inches(3)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = content_box.text_frame
    text_frame.text = "🤖 Browser Agent Assistant\n\n📦 GitHub: BrowserAgent\n🐳 Docker: docker-compose up -d\n📧 联系方式：见个人简历"

    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.space_after = Pt(18)

    return prs

if __name__ == "__main__":
    import sys
    import io

    # 设置UTF-8输出
    if sys.platform == 'win32':
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

    print("Starting PPT generation...")
    prs = create_presentation()

    output_file = "BrowserAgent_Presentation.pptx"
    prs.save(output_file)

    print(f"PPT generated successfully!")
    print(f"File location: {output_file}")
    print(f"Total slides: {len(prs.slides)}")
