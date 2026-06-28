#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
浏览器智能助手 - 专业版PPT生成器
参考Xucheng Zhuang个人简介PPT风格
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 配色方案（参考个人简介PPT）
    NAVY_BLUE = RGBColor(24, 53, 80)      # 深蓝色 #183550
    ACCENT_RED = RGBColor(192, 57, 43)     # 强调红色 #C0392B
    TEXT_DARK = RGBColor(44, 62, 80)       # 深色文字 #2C3E50
    TEXT_LIGHT = RGBColor(149, 165, 166)   # 浅色文字 #95A5A6
    WHITE = RGBColor(255, 255, 255)
    LIGHT_BG = RGBColor(236, 240, 241)     # 浅灰背景 #ECF0F1

    def add_split_slide_title(title, subtitle, author_info=None):
        """创建封面：左侧深蓝，右侧白色"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 左侧深蓝色背景
        left_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(4), Inches(7.5)
        )
        left_bg.fill.solid()
        left_bg.fill.fore_color.rgb = NAVY_BLUE
        left_bg.line.fill.background()

        # 右侧白色背景
        right_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(4), Inches(0),
            Inches(6), Inches(7.5)
        )
        right_bg.fill.solid()
        right_bg.fill.fore_color.rgb = WHITE
        right_bg.line.fill.background()

        # 红色分割条
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(3.9), Inches(0),
            Inches(0.1), Inches(7.5)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = ACCENT_RED
        divider.line.fill.background()

        # 左侧：作者信息
        if author_info:
            # "COMPETITION PRESENTATION"
            label_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(3), Inches(0.4))
            label_frame = label_box.text_frame
            label_frame.text = "COMPETITION PRESENTATION"
            p = label_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_LIGHT
            p.font.name = 'Arial'

            # 项目名称（英文）
            name_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(3), Inches(1.5))
            name_frame = name_box.text_frame
            name_frame.text = "Browser Agent\nAssistant"
            name_frame.word_wrap = True
            for para in name_frame.paragraphs:
                para.font.size = Pt(36)
                para.font.bold = True
                para.font.color.rgb = WHITE
                para.font.name = 'Arial'

            # 项目名称（中文）
            cn_name_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(3), Inches(0.6))
            cn_name_frame = cn_name_box.text_frame
            cn_name_frame.text = "浏览器智能助手"
            p = cn_name_frame.paragraphs[0]
            p.font.size = Pt(20)
            p.font.color.rgb = WHITE
            p.font.name = 'Microsoft YaHei'

            # 红色下划线
            underline = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5), Inches(4),
                Inches(1.2), Inches(0.03)
            )
            underline.fill.solid()
            underline.fill.fore_color.rgb = ACCENT_RED
            underline.line.fill.background()

            # 作者信息
            info_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(3), Inches(2))
            info_frame = info_box.text_frame
            info_frame.text = f"{author_info['role']}\n{author_info['org']}"
            for para in info_frame.paragraphs:
                para.font.size = Pt(14)
                para.font.color.rgb = WHITE
                para.font.name = 'Arial'
                para.space_after = Pt(6)

        # 右侧：标题和副标题
        # 小标签
        right_label_box = slide.shapes.add_textbox(Inches(4.5), Inches(1.5), Inches(5), Inches(0.4))
        right_label_frame = right_label_box.text_frame
        right_label_frame.text = "AI-POWERED WEB AUTOMATION"
        p = right_label_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_LIGHT
        p.font.name = 'Arial'

        # 主标题
        main_title_box = slide.shapes.add_textbox(Inches(4.5), Inches(2.2), Inches(5), Inches(1.5))
        main_title_frame = main_title_box.text_frame
        main_title_frame.text = title
        main_title_frame.word_wrap = True
        for para in main_title_frame.paragraphs:
            para.font.size = Pt(40)
            para.font.bold = True
            para.font.color.rgb = NAVY_BLUE
            para.font.name = 'Arial'
            para.space_after = Pt(12)

        # 红色下划线
        title_underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(4.5), Inches(3.9),
            Inches(1.5), Inches(0.04)
        )
        title_underline.fill.solid()
        title_underline.fill.fore_color.rgb = ACCENT_RED
        title_underline.line.fill.background()

        # 副标题/说明
        subtitle_box = slide.shapes.add_textbox(Inches(4.5), Inches(4.5), Inches(5), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.word_wrap = True
        for para in subtitle_frame.paragraphs:
            para.font.size = Pt(16)
            para.font.color.rgb = TEXT_DARK
            para.font.name = 'Arial'
            para.space_after = Pt(10)

        return slide

    def add_content_slide_with_cards(title, title_en, cards):
        """创建内容页：卡片布局"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 白色背景
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = WHITE
        bg.line.fill.background()

        # 红色顶部条
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.1)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = ACCENT_RED
        top_bar.line.fill.background()

        # 页码（右上角）
        page_num_box = slide.shapes.add_textbox(Inches(9), Inches(0.3), Inches(0.8), Inches(0.3))
        page_num_frame = page_num_box.text_frame
        page_num_frame.text = f"{len(prs.slides)} / 10"
        p = page_num_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_LIGHT
        p.alignment = PP_ALIGN.RIGHT

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.8), Inches(8.8), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = title
        p = title_frame.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = NAVY_BLUE
        p.font.name = 'Microsoft YaHei'

        # 红色下划线
        title_underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.6), Inches(1.5),
            Inches(1.8), Inches(0.03)
        )
        title_underline.fill.solid()
        title_underline.fill.fore_color.rgb = ACCENT_RED
        title_underline.line.fill.background()

        # 卡片内容
        if len(cards) == 2:
            # 2x1布局
            card_width = Inches(4.2)
            card_height = Inches(4.5)
            positions = [(0.6, 2.2), (5.2, 2.2)]
        elif len(cards) == 3:
            # 3列布局
            card_width = Inches(2.8)
            card_height = Inches(4.5)
            positions = [(0.6, 2.2), (3.6, 2.2), (6.6, 2.2)]
        elif len(cards) == 4:
            # 2x2布局
            card_width = Inches(4.2)
            card_height = Inches(2)
            positions = [(0.6, 2.2), (5.2, 2.2), (0.6, 4.5), (5.2, 4.5)]
        else:
            # 默认单列
            card_width = Inches(8.8)
            card_height = Inches(4.5)
            positions = [(0.6, 2.2)]

        for i, card in enumerate(cards):
            if i >= len(positions):
                break

            x, y = positions[i]

            # 卡片边框
            card_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y),
                card_width, card_height
            )
            card_box.fill.solid()
            card_box.fill.fore_color.rgb = LIGHT_BG
            card_box.line.color.rgb = RGBColor(200, 200, 200)
            card_box.line.width = Pt(0.5)

            # 卡片左侧红色条
            if 'label' in card:
                left_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(x), Inches(y),
                    Inches(0.05), card_height
                )
                left_bar.fill.solid()
                left_bar.fill.fore_color.rgb = ACCENT_RED
                left_bar.line.fill.background()

                # 标签
                label_box = slide.shapes.add_textbox(
                    Inches(x + 0.15), Inches(y + 0.15),
                    card_width - Inches(0.3), Inches(0.3)
                )
                label_frame = label_box.text_frame
                label_frame.text = card['label']
                p = label_frame.paragraphs[0]
                p.font.size = Pt(9)
                p.font.color.rgb = ACCENT_RED
                p.font.bold = True
                p.font.name = 'Arial'

            # 卡片标题
            title_y_offset = 0.5 if 'label' in card else 0.2
            card_title_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(y + title_y_offset),
                card_width - Inches(0.4), Inches(0.5)
            )
            card_title_frame = card_title_box.text_frame
            card_title_frame.text = card['title']
            card_title_frame.word_wrap = True
            p = card_title_frame.paragraphs[0]
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = NAVY_BLUE
            p.font.name = 'Microsoft YaHei'

            # 卡片内容
            content_y_offset = 1.1 if 'label' in card else 0.8
            card_content_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(y + content_y_offset),
                card_width - Inches(0.4), card_height - Inches(content_y_offset + 0.2)
            )
            card_content_frame = card_content_box.text_frame
            card_content_frame.word_wrap = True

            for item in card['content']:
                p = card_content_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(13)
                p.font.color.rgb = TEXT_DARK
                p.font.name = 'Microsoft YaHei'
                p.space_after = Pt(8)
                if item.startswith('■'):
                    p.level = 0
                else:
                    p.level = 0

        return slide

    def add_thank_you_slide():
        """创建结束页"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 深蓝色背景
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = NAVY_BLUE
        bg.line.fill.background()

        # 红色顶部条
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(3), Inches(2.5),
            Inches(4), Inches(0.04)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = ACCENT_RED
        top_bar.line.fill.background()

        # Thank You
        title_box = slide.shapes.add_textbox(Inches(2), Inches(2.8), Inches(6), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "Thank You"
        p = title_frame.paragraphs[0]
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Arial'

        # 副标题
        subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(0.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Feel free to ask me any questions."
        p = subtitle_frame.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Arial'

        # 红色底部条
        bottom_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(3), Inches(4.8),
            Inches(4), Inches(0.04)
        )
        bottom_bar.fill.solid()
        bottom_bar.fill.fore_color.rgb = ACCENT_RED
        bottom_bar.line.fill.background()

        # 联系信息
        contact_box = slide.shapes.add_textbox(Inches(2), Inches(5.2), Inches(6), Inches(1.5))
        contact_frame = contact_box.text_frame
        contact_frame.text = "Xucheng Zhuang  ·  Sichuan University\nscuzxc660223@gmail.com\ngithub.com/Tonyzxch"
        for para in contact_frame.paragraphs:
            para.font.size = Pt(14)
            para.font.color.rgb = WHITE
            para.alignment = PP_ALIGN.CENTER
            para.font.name = 'Arial'
            para.space_after = Pt(8)

        return slide

    # ===== 开始创建幻灯片 =====

    # 第1页：封面
    add_split_slide_title(
        "Autonomous Browser\nOperation Assistant",
        "RESEARCH FOCUS\n\nNatural Language Understanding  ·  Web Automation  ·  AI Agent Systems\n\n\nPROJECT OVERVIEW\n\nAn AI-powered Chrome extension that enables users to automate complex web tasks through natural language commands.",
        author_info={
            'role': 'B.E. in Software Engineering',
            'org': 'Sichuan University'
        }
    )

    # 第2页：项目背景
    add_content_slide_with_cards(
        "项目背景 Background",
        "",
        [
            {
                'label': 'TASK OBJECTIVE',
                'title': '开发智能浏览器插件',
                'content': [
                    '通过自然语言实现网页自动化',
                    '让AI成为用户的浏览器助手',
                    '降低自动化操作的技术门槛'
                ]
            },
            {
                'label': 'CORE CHALLENGES',
                'title': '技术挑战',
                'content': [
                    '■ 理解用户意图和任务目标',
                    '■ 自动感知页面结构',
                    '■ 智能规划执行步骤',
                    '■ 准确定位并操作元素'
                ]
            }
        ]
    )

    # 第3页：核心功能（4个卡片）
    add_content_slide_with_cards(
        "核心功能 Key Features",
        "",
        [
            {
                'title': '🤖 自然语言交互',
                'content': [
                    '用中文描述任务',
                    'AI自动理解执行'
                ]
            },
            {
                'title': '🔍 智能页面感知',
                'content': [
                    '自动识别可交互元素',
                    '生成精准选择器'
                ]
            },
            {
                'title': '🧠 AI任务规划',
                'content': [
                    'Claude Opus 4.8驱动',
                    '多步骤任务分解'
                ]
            },
            {
                'title': '⚡ 自动化执行',
                'content': [
                    '点击、输入、导航',
                    '信息提取等操作'
                ]
            }
        ]
    )

    # 第4页：技术架构
    add_content_slide_with_cards(
        "技术架构 Technical Stack",
        "",
        [
            {
                'label': 'FRONTEND',
                'title': 'React + TypeScript',
                'content': [
                    'React 18框架',
                    'TypeScript类型安全',
                    'Vite 5构建工具',
                    'Lucide图标库'
                ]
            },
            {
                'label': 'EXTENSION',
                'title': 'Chrome Extension',
                'content': [
                    'Manifest V3标准',
                    'Background Service Worker',
                    'Content Script注入',
                    'Chrome Messaging API'
                ]
            },
            {
                'label': 'AI SERVICE',
                'title': 'Claude Opus 4.8',
                'content': [
                    'OpenAI兼容API',
                    '结构化Prompt工程',
                    'JSON格式化响应',
                    '智能任务规划'
                ]
            },
            {
                'label': 'DEPLOYMENT',
                'title': 'Docker + Nginx',
                'content': [
                    'Docker容器化',
                    'Nginx静态服务',
                    '一键部署',
                    '生产就绪'
                ]
            }
        ]
    )

    # 第5页：设计亮点
    add_content_slide_with_cards(
        "设计亮点 Design Highlights",
        "",
        [
            {
                'title': '01 · 智能元素定位',
                'content': [
                    '■ CSS Selector + XPath双重策略',
                    '■ 自动过滤隐藏元素',
                    '■ 支持文本内容匹配',
                    '■ 限制选择器深度避免脆弱'
                ]
            },
            {
                'title': '02 · AI驱动规划',
                'content': [
                    '■ 精心设计的Prompt工程',
                    '■ 结构化JSON响应',
                    '■ 任务步骤智能分解',
                    '■ 上下文感知决策'
                ]
            },
            {
                'title': '03 · 实时状态同步',
                'content': [
                    '■ Chrome Messaging广播',
                    '■ 多组件响应式更新',
                    '■ 步骤级粒度追踪',
                    '■ UI自动刷新'
                ]
            },
            {
                'title': '04 · 模块化架构',
                'content': [
                    '■ 清晰的职责分离',
                    '■ TypeScript类型安全',
                    '■ 易于扩展维护',
                    '■ 代码复用性高'
                ]
            }
        ]
    )

    # 第6页：实现成果
    add_content_slide_with_cards(
        "实现成果 Achievements",
        "",
        [
            {
                'label': 'FUNCTIONALITY',
                'title': '功能完整性',
                'content': [
                    '任务理解、页面感知、规划执行、结果反馈完整闭环',
                    '支持7种动作类型：navigate、click、input、extract、scroll、wait、analyze'
                ]
            },
            {
                'label': 'CODE QUALITY',
                'title': '代码质量',
                'content': [
                    '2000+行代码，TypeScript类型安全',
                    '11个源码文件，模块化设计清晰',
                    '完善的错误处理和日志记录'
                ]
            },
            {
                'label': 'DOCUMENTATION',
                'title': '文档完整',
                'content': [
                    '11个Markdown文档',
                    '涵盖安装、使用、开发全流程',
                    '详细的API配置和调试指南'
                ]
            },
            {
                'label': 'DEPLOYMENT',
                'title': '生产就绪',
                'content': [
                    'Docker一键部署',
                    '插件大小仅65KB',
                    '构建时间<3秒'
                ]
            }
        ]
    )

    # 第7页：使用示例
    add_content_slide_with_cards(
        "使用示例 Use Cases",
        "",
        [
            {
                'title': '场景1：GitHub项目搜索',
                'content': [
                    '"帮我在GitHub搜索React UI组件库，按star排序"',
                    '',
                    '→ 自动定位搜索框',
                    '→ 输入"React UI library"',
                    '→ 点击搜索按钮',
                    '→ 应用排序筛选'
                ]
            },
            {
                'title': '场景2：豆瓣电影查询',
                'content': [
                    '"在豆瓣搜索评分8分以上的科幻电影"',
                    '',
                    '→ 导航到豆瓣电影',
                    '→ 输入搜索关键词',
                    '→ 设置评分筛选',
                    '→ 提取电影信息'
                ]
            },
            {
                'title': '场景3：电商商品筛选',
                'content': [
                    '"在淘宝搜索500元以内的蓝牙耳机"',
                    '',
                    '→ 搜索商品',
                    '→ 设置价格区间',
                    '→ 提取商品列表',
                    '→ 生成对比表格'
                ]
            }
        ]
    )

    # 第8页：开发过程
    add_content_slide_with_cards(
        "开发过程 Development Process",
        "",
        [
            {
                'label': 'DEVELOPMENT TIME',
                'title': '2小时完成开发',
                'content': [
                    '使用Claude Code辅助开发',
                    '一次性完成核心功能实现',
                    '6次Git提交，清晰的版本历史'
                ]
            },
            {
                'label': 'ITERATION STEPS',
                'title': '迭代过程',
                'content': [
                    '1️⃣ 需求分析与架构设计（10分钟）',
                    '2️⃣ 核心功能开发（60分钟）',
                    '3️⃣ UI界面与交互优化（30分钟）',
                    '4️⃣ Docker部署配置（15分钟）',
                    '5️⃣ 文档编写与完善（15分钟）'
                ]
            }
        ]
    )

    # 第9页：未来规划
    add_content_slide_with_cards(
        "未来规划 Future Plans",
        "",
        [
            {
                'title': '01 · 功能增强',
                'content': [
                    'Cross-platform Support',
                    '支持更多动作类型（拖拽、右键菜单）',
                    '添加任务模板系统',
                    '多标签页协同操作'
                ]
            },
            {
                'title': '02 · 性能优化',
                'content': [
                    'Performance Optimization',
                    '减少页面信息传输量',
                    '缓存AI规划结果',
                    '并行执行独立步骤'
                ]
            },
            {
                'title': '03 · 平台扩展',
                'content': [
                    'Platform Extension',
                    '支持Firefox等浏览器',
                    '添加本地模型选项',
                    '开放API供第三方集成'
                ]
            }
        ]
    )

    # 第10页：结束页
    add_thank_you_slide()

    return prs

if __name__ == "__main__":
    import sys
    import io

    # 设置UTF-8输出
    if sys.platform == 'win32':
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

    print("Generating professional PPT...")
    prs = create_presentation()

    output_file = "BrowserAgent_Presentation_Professional.pptx"
    prs.save(output_file)

    print(f"PPT generated successfully!")
    print(f"File: {output_file}")
    print(f"Slides: {len(prs.slides)}")
